#!/usr/bin/env python3
import csv
import io
import json
import mimetypes
import os
import sys
import zipfile
import tarfile
from pathlib import Path
from typing import Any, Dict, List, Optional

SENSITIVE_NAMES = {
    ".env",
    ".env.local",
    ".env.production",
    "id_rsa",
    "id_ed25519",
    "known_hosts",
    "credentials",
    "credentials.json",
    "token",
    "token.json",
}

SENSITIVE_SUFFIXES = {
    ".pem",
    ".key",
    ".p12",
    ".pfx",
}

TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".markdown",
    ".json",
    ".jsonl",
    ".yaml",
    ".yml",
    ".xml",
    ".html",
    ".htm",
    ".csv",
    ".tsv",
    ".log",
    ".ini",
    ".toml",
}


def configure_tesseract_windows() -> None:
    if os.name != "nt":
        return
    candidates = [
        Path(r"C:\Program Files\Tesseract-OCR\tesseract.exe"),
        Path(r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"),
    ]
    for candidate in candidates:
        if candidate.exists():
            os.environ["PATH"] = str(candidate.parent) + os.pathsep + os.environ.get("PATH", "")
            try:
                import pytesseract  # type: ignore
                pytesseract.pytesseract.tesseract_cmd = str(candidate)
            except Exception:
                pass
            return


def fail(message: str, code: int = 1) -> None:
    print(json.dumps({"ok": False, "error": message}, ensure_ascii=False))
    raise SystemExit(code)


def safe_read_text(path: Path, max_bytes: int = 8_000_000) -> str:
    data = path.read_bytes()[:max_bytes]
    for enc in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return data.decode(enc, errors="replace")
        except Exception:
            pass
    return data.decode("utf-8", errors="replace")


def check_sensitive(path: Path) -> None:
    lower_name = path.name.lower()
    if lower_name in SENSITIVE_NAMES:
        fail(f"Refusing to read sensitive file: {path.name}")
    if path.suffix.lower() in SENSITIVE_SUFFIXES:
        fail(f"Refusing to read sensitive key/certificate-like file: {path.name}")
    parts = {p.lower() for p in path.parts}
    if ".ssh" in parts or ".gnupg" in parts:
        fail(f"Refusing to read file under sensitive directory: {path}")


def truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + f"\n\n[TRUNCATED: original_chars={len(text)}, max_chars={max_chars}]"


def make_record(
    path: Path,
    kind: str,
    content: str,
    metadata: Optional[Dict[str, Any]] = None,
    chunks: Optional[List[Dict[str, Any]]] = None,
) -> Dict[str, Any]:
    stat = path.stat()
    mime, _ = mimetypes.guess_type(str(path))
    return {
        "ok": True,
        "path": str(path),
        "name": path.name,
        "suffix": path.suffix.lower(),
        "kind": kind,
        "size_bytes": stat.st_size,
        "mime_guess": mime,
        "metadata": metadata or {},
        "chunks": chunks or [],
        "content": content,
    }


def read_pdf(path: Path, page_start: Optional[int], page_end: Optional[int]) -> Dict[str, Any]:
    try:
        import pypdf
    except Exception as e:
        fail(f"Missing dependency for PDF: pypdf. Install requirements-docs.txt. Detail: {e}")

    reader = pypdf.PdfReader(str(path))
    total = len(reader.pages)

    start = 1 if page_start is None else max(1, int(page_start))
    end = total if page_end is None else min(total, int(page_end))
    if start > end:
        fail(f"Invalid page range: page_start={start}, page_end={end}, total_pages={total}")

    chunks: List[Dict[str, Any]] = []
    all_text: List[str] = []

    for page_num in range(start, end + 1):
        page = reader.pages[page_num - 1]
        try:
            text = page.extract_text() or ""
        except Exception as e:
            text = f"[PDF text extraction error on page {page_num}: {e}]"
        chunks.append({
            "type": "page",
            "page": page_num,
            "text": text,
        })
        all_text.append(f"\n\n## Page {page_num}\n\n{text}")

    metadata = {
        "total_pages": total,
        "page_start": start,
        "page_end": end,
        "reader": "pypdf",
    }

    if reader.metadata:
        metadata["pdf_metadata"] = {str(k): str(v) for k, v in reader.metadata.items()}

    return make_record(path, "pdf", "\n".join(all_text).strip(), metadata, chunks)


def read_docx(path: Path) -> Dict[str, Any]:
    try:
        import docx
    except Exception as e:
        fail(f"Missing dependency for DOCX: python-docx. Detail: {e}")

    doc = docx.Document(str(path))
    parts: List[str] = []
    chunks: List[Dict[str, Any]] = []

    for i, p in enumerate(doc.paragraphs):
        text = p.text.strip()
        if text:
            chunks.append({"type": "paragraph", "index": i, "text": text})
            parts.append(text)

    for ti, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        chunks.append({"type": "table", "index": ti, "rows": rows})
        parts.append(f"\n\n### Table {ti + 1}\n" + "\n".join([" | ".join(r) for r in rows]))

    return make_record(path, "docx", "\n\n".join(parts).strip(), {"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}, chunks)


def read_xlsx(path: Path) -> Dict[str, Any]:
    try:
        import openpyxl
    except Exception as e:
        fail(f"Missing dependency for XLSX: openpyxl. Detail: {e}")

    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    parts: List[str] = []
    chunks: List[Dict[str, Any]] = []

    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            values = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in values):
                rows.append(values)

        preview = rows[:200]
        chunks.append({
            "type": "sheet",
            "sheet": ws.title,
            "rows_preview": preview,
            "row_count_previewed": len(preview),
        })

        parts.append(f"\n\n## Sheet: {ws.title}\n")
        for r in preview:
            parts.append(" | ".join(r))

    return make_record(path, "xlsx", "\n".join(parts).strip(), {"sheets": wb.sheetnames}, chunks)


def read_pptx(path: Path) -> Dict[str, Any]:
    try:
        from pptx import Presentation
    except Exception as e:
        fail(f"Missing dependency for PPTX: python-pptx. Detail: {e}")

    prs = Presentation(str(path))
    parts: List[str] = []
    chunks: List[Dict[str, Any]] = []

    for si, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
        text = "\n".join([t for t in texts if t])
        chunks.append({"type": "slide", "slide": si, "text": text})
        parts.append(f"\n\n## Slide {si}\n{text}")

    return make_record(path, "pptx", "\n".join(parts).strip(), {"slides": len(prs.slides)}, chunks)


def read_csv_like(path: Path) -> Dict[str, Any]:
    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
    text = safe_read_text(path)
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows = []
    for i, row in enumerate(reader):
        if i >= 500:
            break
        rows.append(row)

    content = "\n".join([" | ".join(r) for r in rows])
    return make_record(path, path.suffix.lower().lstrip("."), content, {"rows_previewed": len(rows)}, [{"type": "table", "rows": rows}])


def read_html(path: Path) -> Dict[str, Any]:
    raw = safe_read_text(path)
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(raw, "html.parser")
        title = soup.title.string.strip() if soup.title and soup.title.string else ""
        text = soup.get_text("\n", strip=True)
        return make_record(path, "html", text, {"title": title})
    except Exception:
        return make_record(path, "html", raw, {"parser": "raw-fallback"})


def read_image(path: Path, ocr: bool) -> Dict[str, Any]:
    try:
        from PIL import Image
    except Exception as e:
        fail(f"Missing dependency for image reading: Pillow. Detail: {e}")

    img = Image.open(str(path))
    metadata = {
        "format": img.format,
        "mode": img.mode,
        "width": img.width,
        "height": img.height,
    }

    text = f"Image file: {path.name}\nFormat: {img.format}\nSize: {img.width}x{img.height}\nMode: {img.mode}"

    if ocr:
        try:
            configure_tesseract_windows()
            import pytesseract
            ocr_text = pytesseract.image_to_string(img)
            text += "\n\n## OCR Text\n\n" + ocr_text
            metadata["ocr"] = "pytesseract"
        except Exception as e:
            text += f"\n\n[OCR unavailable or failed: {e}]"
            metadata["ocr_error"] = str(e)

    return make_record(path, "image", text, metadata)


def read_archive(path: Path) -> Dict[str, Any]:
    entries = []
    if zipfile.is_zipfile(path):
        with zipfile.ZipFile(path) as zf:
            for info in zf.infolist()[:500]:
                entries.append({
                    "name": info.filename,
                    "size": info.file_size,
                    "compressed_size": info.compress_size,
                })
        content = "\n".join([f"{e['name']} | {e['size']} bytes" for e in entries])
        return make_record(path, "zip", content, {"entries_previewed": len(entries)}, [{"type": "archive_entries", "entries": entries}])

    if tarfile.is_tarfile(path):
        with tarfile.open(path) as tf:
            for m in tf.getmembers()[:500]:
                entries.append({"name": m.name, "size": m.size, "type": "dir" if m.isdir() else "file"})
        content = "\n".join([f"{e['name']} | {e['size']} bytes | {e['type']}" for e in entries])
        return make_record(path, "tar", content, {"entries_previewed": len(entries)}, [{"type": "archive_entries", "entries": entries}])

    fail(f"Unsupported archive format: {path}")


def read_epub(path: Path) -> Dict[str, Any]:
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except Exception as e:
        fail(f"Missing dependency for EPUB: ebooklib and beautifulsoup4. Detail: {e}")

    book = epub.read_epub(str(path))
    parts = []
    chunks = []

    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            html = item.get_content().decode("utf-8", errors="replace")
            soup = BeautifulSoup(html, "html.parser")
            text = soup.get_text("\n", strip=True)
            if text:
                chunks.append({"type": "epub_document", "name": item.get_name(), "text": text})
                parts.append(f"\n\n## {item.get_name()}\n{text}")

    return make_record(path, "epub", "\n".join(parts).strip(), {"documents": len(chunks)}, chunks)


def dispatch(args: Dict[str, Any]) -> Dict[str, Any]:
    target = args.get("target")
    if not target:
        fail("Missing required field: target")

    path = Path(target).expanduser().resolve()
    if not path.exists():
        fail(f"File not found: {path}")
    if not path.is_file():
        fail(f"Target is not a file: {path}")

    check_sensitive(path)

    suffix = path.suffix.lower()
    page_start = args.get("pageStart") or args.get("page_start")
    page_end = args.get("pageEnd") or args.get("page_end")
    ocr = bool(args.get("ocr", False))

    if suffix == ".pdf":
        return read_pdf(path, page_start, page_end)
    if suffix == ".docx":
        return read_docx(path)
    if suffix == ".xlsx":
        return read_xlsx(path)
    if suffix == ".pptx":
        return read_pptx(path)
    if suffix in {".csv", ".tsv"}:
        return read_csv_like(path)
    if suffix in {".html", ".htm"}:
        return read_html(path)
    if suffix == ".epub":
        return read_epub(path)
    if suffix in {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp"}:
        return read_image(path, ocr)
    if suffix in {".zip", ".tar", ".gz", ".tgz"} or zipfile.is_zipfile(path) or tarfile.is_tarfile(path):
        return read_archive(path)
    if suffix in TEXT_SUFFIXES or path.stat().st_size < 4_000_000:
        text = safe_read_text(path)
        return make_record(path, "text", text)

    fail(f"Unsupported file type: {suffix}")


def to_markdown(record: Dict[str, Any], max_chars: int) -> str:
    if not record.get("ok"):
        return json.dumps(record, ensure_ascii=False, indent=2)

    meta = record.get("metadata") or {}
    lines = [
        f"# Document extraction: {record.get('name')}",
        "",
        f"- kind: {record.get('kind')}",
        f"- path: {record.get('path')}",
        f"- suffix: {record.get('suffix')}",
        f"- size_bytes: {record.get('size_bytes')}",
        f"- mime_guess: {record.get('mime_guess')}",
        "",
        "## Metadata",
        "",
        "```json",
        json.dumps(meta, ensure_ascii=False, indent=2),
        "```",
        "",
        "## Content",
        "",
        record.get("content") or "",
    ]
    return truncate("\n".join(lines), max_chars)


def main() -> None:
    try:
        raw = sys.stdin.read()
        args = json.loads(raw) if raw.strip() else {}
    except Exception as e:
        fail(f"Invalid JSON input: {e}")

    max_chars = int(args.get("maxChars") or args.get("max_chars") or 20000)
    include_json = bool(args.get("includeJson") or args.get("include_json") or False)

    record = dispatch(args)

    if include_json:
        record["content"] = truncate(record.get("content") or "", max_chars)
        for chunk in record.get("chunks", []):
            if isinstance(chunk.get("text"), str):
                chunk["text"] = truncate(chunk["text"], max_chars)
        print(json.dumps(record, ensure_ascii=False, indent=2))
    else:
        print(to_markdown(record, max_chars))


if __name__ == "__main__":
    main()
